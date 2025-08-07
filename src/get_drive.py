#!/usr/bin/env python
# -*- coding: utf-8 -*-

import json
from io import BytesIO
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

SCOPES = [
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/spreadsheets.readonly'
]
SERVICE_ACCOUNT_FILE = 'service_account/slack-ai-52557-df876200708f.json'  # 適宜変更してください

creds = service_account.Credentials.from_service_account_file(
    SERVICE_ACCOUNT_FILE, scopes=SCOPES)
drive_service = build('drive', 'v3', credentials=creds)

def get_file_content(file_id, mimeType):
    content = ""
    try:
        if mimeType == "application/vnd.google-apps.document":
            print(f"Document {file_id} をテキスト形式でエクスポート中...")
            request = drive_service.files().export_media(fileId=file_id, mimeType="text/plain")
            fh = BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                print(f"Downloading {file_id}: {int(status.progress() * 100)}%")
            content = fh.getvalue().decode("utf-8")
        elif mimeType == "application/vnd.google-apps.spreadsheet":
            print(f"Spreadsheet {file_id} の全シートのデータを取得中...")
            sheets_service = build('sheets', 'v4', credentials=creds)
            spreadsheet = sheets_service.spreadsheets().get(
                spreadsheetId=file_id,
                includeGridData=True
            ).execute()
            all_sheet_texts = []
            for sheet in spreadsheet.get('sheets', []):
                sheet_title = sheet.get('properties', {}).get('title', 'Sheet')
                sheet_text = f"シート: {sheet_title}\n"
                grid_data = sheet.get('data', [])
                for grid in grid_data:
                    row_data = grid.get('rowData', [])
                    for row in row_data:
                        cell_values = []
                        for cell in row.get('values', []):
                            cell_text = cell.get('formattedValue', '')
                            cell_values.append(cell_text)
                        sheet_text += "\t".join(cell_values) + "\n"
                all_sheet_texts.append(sheet_text)
            content = "\n".join(all_sheet_texts)
        elif mimeType == "application/vnd.google-apps.presentation":
            print(f"Presentation {file_id} を PDF 形式でエクスポート中...")
            request = drive_service.files().export_media(fileId=file_id, mimeType="application/pdf")
            fh = BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                print(f"Downloading {file_id}: {int(status.progress() * 100)}%")
            import PyPDF2
            fh.seek(0)
            reader = PyPDF2.PdfReader(fh)
            texts = []
            for page in reader.pages:
                texts.append(page.extract_text())
            content = "\n".join(texts)
        elif mimeType == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            print(f"DOCX {file_id} をダウンロード中...")
            request = drive_service.files().get_media(fileId=file_id)
            fh = BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                print(f"Downloading {file_id}: {int(status.progress() * 100)}%")
            from docx import Document
            fh.seek(0)
            document = Document(fh)
            content = "\n".join([para.text for para in document.paragraphs])
        elif mimeType == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            print(f"PPTX {file_id} をダウンロード中...")
            request = drive_service.files().get_media(fileId=file_id)
            fh = BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                print(f"Downloading {file_id}: {int(status.progress() * 100)}%")
            from pptx import Presentation
            fh.seek(0)
            prs = Presentation(fh)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            content = "\n".join(texts)
        elif mimeType == "application/pdf":
            print(f"PDF {file_id} をダウンロード中...")
            request = drive_service.files().get_media(fileId=file_id)
            fh = BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                print(f"Downloading {file_id}: {int(status.progress() * 100)}%")
            import PyPDF2
            fh.seek(0)
            reader = PyPDF2.PdfReader(fh)
            texts = []
            for page in reader.pages:
                texts.append(page.extract_text())
            content = "\n".join(texts)
        elif mimeType == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            print(f"XLSX {file_id} をダウンロード中...")
            request = drive_service.files().get_media(fileId=file_id)
            fh = BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                print(f"Downloading {file_id}: {int(status.progress() * 100)}%")
            from openpyxl import load_workbook
            fh.seek(0)
            wb = load_workbook(fh, read_only=True, data_only=True)
            sheet_texts = []
            for ws in wb.worksheets:
                sheet_text = f"シート: {ws.title}\n"
                for row in ws.iter_rows(values_only=True):
                    row_str = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    sheet_text += row_str + "\n"
                sheet_texts.append(sheet_text)
            content = "\n".join(sheet_texts)
        else:
            print(f"ファイル {file_id} は対応していない MIMEタイプ: {mimeType}")
    except Exception as e:
        print(f"ファイル {file_id} のコンテンツ取得に失敗: {e}")
    return content

def list_files_recursive(folder_id, files_data=None):
    if files_data is None:
        files_data = []
    page_token = None
    query = f"'{folder_id}' in parents"
    while True:
        try:
            response = drive_service.files().list(
                q=query,
                fields="nextPageToken, files(id, name, mimeType, webViewLink, createdTime, modifiedTime, owners)",
                pageToken=page_token,
                pageSize=1000
            ).execute()
        except Exception as e:
            print(f"フォルダ {folder_id} のファイルリスト取得に失敗: {e}")
            break

        for file in response.get('files', []):
            if file.get("mimeType") == "application/vnd.google-apps.folder":
                try:
                    list_files_recursive(file.get("id"), files_data)
                except Exception as e:
                    print(f"フォルダ {file.get('id')} の処理中にエラー: {e}")
            else:
                try:
                    perms = drive_service.permissions().list(
                        fileId=file.get("id"),
                        fields="permissions(id, type, role, emailAddress, displayName, domain)"
                    ).execute()
                except Exception as e:
                    print(f"ファイル {file.get('id')} のアクセス権取得に失敗: {e}")
                    perms = {"permissions": []}

                has_domain_permission = any(perm.get("type") == "domain" and perm.get("domain", "") == "techfund.jp"
                                            for perm in perms.get("permissions", []))
                if has_domain_permission:
                    collaborators = []
                    for perm in perms.get("permissions", []):
                        if perm.get("type") == "user" and perm.get("role") in ["writer", "commenter"]:
                            collaborators.append({
                                "id": perm.get("id"),
                                "displayName": perm.get("displayName"),
                                "emailAddress": perm.get("emailAddress"),
                                "role": perm.get("role")
                            })
                    file_info = {
                        "id": file.get("id"),
                        "title": file.get("name"),
                        "url": file.get("webViewLink") if file.get("webViewLink") else "",
                        "createdTime": file.get("createdTime"),
                        "modifiedTime": file.get("modifiedTime"),
                        "owners": file.get("owners"),
                        "collaborators": collaborators,
                        "content": get_file_content(file.get("id"), file.get("mimeType"))
                    }
                    files_data.append(file_info)
        page_token = response.get("nextPageToken", None)
        if not page_token:
            break
    return files_data

if __name__ == '__main__':
    folder_id = "179ksE67kVo3PXEZbWJRj2zcg0qUAoWsm"  # TECHFUND Inc.フォルダのID
    data = list_files_recursive(folder_id)
    with open("drive_documents.json", "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)